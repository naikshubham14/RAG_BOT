import os
import sys
sys.path.append("../..")
from dotenv import load_dotenv, find_dotenv
from PyPDF2 import PdfReader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.embeddings import HuggingFaceInferenceAPIEmbeddings
from langchain.vectorstores import Chroma, FAISS
from langchain.llms import HuggingFaceHub
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain
from pptx import Presentation

_ = load_dotenv(find_dotenv())
HuggingFaceInferenceAPIKey = os.environ['HUGGINGFACEHUB_API_TOKEN']

def read_pdf(pdf):
    """
    This function reads a PDF file and extracts its text content.

    Parameters:
    pdf (str): The path to the PDF file.

    Returns:
    str: The extracted text from the PDF file.
    """
    text = ""
    reader = PdfReader(pdf)
    for page in reader.pages:
        text += page.extract_text()
    return text

def split_text(text):
    """
    This function splits a large text into smaller chunks using the RecursiveCharacterTextSplitter.
    The splitter divides the text into chunks of a specified size, with an overlap of characters between chunks.
    The chunks are separated by a newline character.

    Parameters:
    text (str): The large text to be split.

    Returns:
    List[str]: A list of smaller text chunks.
    """
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200,
        separators=["\n\n", "\n", " ", ""],
        length_function=len)
    return text_splitter.split_text(text)

def load_vector_db(chunks):
    """
    This function initializes and loads a vector database (Chroma) with text chunks.
    The function removes any existing Chroma directory, creates a new one, and initializes
    a Chroma vector store with the provided text chunks.

    Parameters:
    chunks (List[str]): A list of text chunks to be loaded into the vector database.

    Returns:
    Chroma: An initialized Chroma vector store with the provided text chunks.
    """
    """
    directory = 'chroma/'
    subprocess.run(['rm', '-rf', directory], check=True)
    embeddings = HuggingFaceInferenceAPIEmbeddings(
        api_key=HuggingFaceInferenceAPIKey,
        model_name="sentence-transformers/all-MiniLM-l6-v2")
    vectorstore = Chroma.from_documents(embeddings, chunks, persist_directory=directory)"""

    embeddings = HuggingFaceInferenceAPIEmbeddings(
        api_key=HuggingFaceInferenceAPIKey,
        model_name="sentence-transformers/all-MiniLM-l6-v2")
    vectordb = FAISS.from_texts(texts=chunks, embedding=embeddings)

    return vectordb

def init_conversation_chain(vectordb):
    """
    Initializes a ConversationalRetrievalChain for a given vector database.

    The ConversationalRetrievalChain is a chain that uses a language model (LLM) and a retriever
    to generate responses based on a given input. The LLM is used to generate a response,
    and the retriever is used to find relevant documents or information in the vector database.

    Parameters:
    vectordb (Chroma or FAISS): The vector database to be used for information retrieval.

    Returns:
    ConversationalRetrievalChain: An initialized ConversationalRetrievalChain object.
    """
    llm = HuggingFaceHub(repo_id="mistralai/Mistral-7B-Instruct-v0.2", model_kwargs={"temperature":0.01, "max_length":4000})
    memory = ConversationBufferMemory(memory_key="chat_history", return_messages=True)
    chain = ConversationalRetrievalChain.from_llm(llm=llm, retriever=vectordb.as_retriever(), memory=memory)
    return chain

def retrieve_build_context(vector_db, user_query):
    """
    Retrieves and builds a contextual background for a given user query using a vector database.

    This function performs a maximal marginal relevance search on the vector database to find the most relevant documents
    based on the user query. It then extracts the content of these documents and concatenates them to form a contextual
    background.

    Parameters:
    vector_db (Chroma or FAISS): The vector database to be used for information retrieval. The database should be
        initialized with documents containing the 'page_content' attribute.
    user_query (str): The user's query for which the contextual background needs to be retrieved.

    Returns:
    str: A concatenated string containing the content of the most relevant documents found in the vector database.
    """
    results = vector_db.max_marginal_relevance_search(user_query)
    context = ""
    for result in results:
        context += result.page_content
    return context

def build_query(context, user_query):
    """
    This function constructs a query string for a language model (LLM) based on a given context and user query.
    The constructed query includes a base prompt, the provided context, and the user's question.
    The base prompt is formatted to guide the LLM in understanding the context and answering the question.

    Parameters:
    context (str): The contextual background information that the LLM should consider when answering the question.
        The context should be a string containing the relevant information.
    user_query (str): The question that the LLM needs to answer. The question should be a string.

    Returns:
    str: A query string that can be used as input to a language model. The query string includes the base prompt,
        the provided context, and the user's question.
    """

    base_prompt = '''Use the following pieces of context delimeted by triple ` to answer the question at the end.
    Provide the answer only if its preset in the given context. Keep you answer as concise as possible.
    If you don't know the answer, just say that you don't know, don't try to make up an answer.

    ```
    context: '''

    final_query =  base_prompt + context + '''\n```\nquestion:'''+ user_query+ '\nAnswer:'
    return final_query

def generate_response(llm, vector_db, user_query):
    """
    This function generates a response based on a given language model (LLM) and query.
    The function takes an LLM and a query as input, sends the query to the LLM,
    and extracts the response from the LLM's output.

    Parameters:
    llm (HuggingFaceHub): The language model to be used for generating the response.
        The LLM should be initialized with a suitable model and parameters.
    query (str): The query to be sent to the LLM. The query should be a string
        containing the user's question or instruction.

    Returns:
    str: The response generated by the LLM. The response should be a string
        containing the LLM's answer to the user's question or fulfillment of the user's instruction.
    """
    context = retrieve_build_context(vector_db, user_query)
    final_query = build_query(context, user_query)
    response = llm(final_query, verbose=True)
    return response.split('Answer:')[1]

def init_llm():
    """
    Initializes and returns a language model (LLM) using the HuggingFaceHub.

    The function initializes an LLM with a specified model repository ID and model parameters.
    The model repository ID is "mistralai/Mistral-7B-Instruct-v0.2", and the model parameters include
    temperature (0.001) and max_length (4000).

    Parameters:
    None

    Returns:
    HuggingFaceHub: An initialized LLM object using the specified model repository ID and parameters.
    """
    llm = HuggingFaceHub(repo_id="mistralai/Mistral-7B-Instruct-v0.2", model_kwargs={"temperature":0.001, "max_length":32000})
    return llm

def parse_file_to_text(file):
    """
+    This function reads a file and extracts its text content based on the file extension.
+    It supports PDF, PowerPoint (PPT/PPTX), and plain text (TXT) files.
+
+    Parameters:
+    file (str): The path to the file to be read. The file must be a PDF, PowerPoint, or TXT file.
+
+    Returns:
+    str: The extracted text content from the file. If the file extension is not supported, an empty string is returned.
+    """
    _, extension = os.path.splitext(file)
    text = ""
    if extension == ".pdf":
        reader = PdfReader(file)
        for page in reader.pages:
            text += page.extract_text()    
    if extension == ".ppt" or extension == ".pptx":
        prs = Presentation(file)
        for slide_number, slide in enumerate(prs.slides):
            for shape in slide.shapes: 
                if hasattr(shape, "text"): 
                    text += shape.text
    if extension == ".txt":
        with open(file, "r") as file:
            text = file.read()    
    return text

        

